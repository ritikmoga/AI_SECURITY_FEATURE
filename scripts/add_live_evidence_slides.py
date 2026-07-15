"""Append live scanner test evidence slides to the ScamShield PowerPoint."""
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.oxml.xmlchemy import OxmlElement


ROOT = Path(__file__).resolve().parents[1]
DECK = ROOT / "docs" / "ScamShield_AI_Project_Presentation.pptx"
SHOTS = ROOT / "docs" / "live-screenshots"
BG, PANEL, RED, PINK, WHITE, MUTED = "080304", "1A070A", "E11D3A", "FF8A99", "FFF5F5", "BFA6A9"


def rgb(value): return RGBColor.from_string(value)


def text(slide, x, y, w, h, value, size=16, shade=WHITE, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h)); frame = box.text_frame; frame.clear(); frame.word_wrap = True
    para = frame.paragraphs[0]; para.alignment = align; run = para.add_run(); run.text = value; run.font.name = "Aptos"; run.font.size = Pt(size); run.font.bold = bold; run.font.color.rgb = rgb(shade)


def box(slide, x, y, label, emphasis=False):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(1.5), Inches(.72)); shape.fill.solid(); shape.fill.fore_color.rgb = rgb(RED if emphasis else PANEL); shape.line.color.rgb = rgb(RED); shape.line.transparency = 20
    text(slide, x+.08, y+.17, 1.34, .3, label, 11, WHITE, True, PP_ALIGN.CENTER)


def arrow(slide, x1, y1, x2, y2):
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2)); line.line.color.rgb = rgb(RED); line.line.width = Pt(1.6); line.line.end_arrowhead = True


def transition(slide):
    item = OxmlElement("p:transition"); item.set("spd", "slow"); fade = OxmlElement("p:fade"); item.append(fade); slide._element.insert(-1, item)


def add(slide, number, title, input_text, detail_file, verdict):
    slide.background.fill.solid(); slide.background.fill.fore_color.rgb = rgb(BG); transition(slide)
    top = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(.08)); top.fill.solid(); top.fill.fore_color.rgb = rgb(RED); top.line.fill.background()
    text(slide, .55, .45, 12, .3, "LIVE TEST EVIDENCE", 10, RED, True); text(slide, .55, .8, 12, .56, title, 27, WHITE, True); text(slide, 11.7, .45, 1, .25, f"{number:02d}", 9, MUTED, True, PP_ALIGN.RIGHT)
    text(slide, .6, 1.55, 5.35, .24, "SAFE TEST INPUT", 9, RED, True); text(slide, .6, 1.83, 5.35, .78, input_text, 15, PINK, False)
    # Diagram connectors first.
    for x in (1.7, 3.45, 5.2): arrow(slide, x, 3.82, x+.7, 3.82)
    for x, label in zip((.5, 2.25, 4.0, 5.75), ("INPUT", "API", "SCORER", "REPORT")): box(slide, x, 3.45, label, label == "REPORT")
    text(slide, .6, 4.65, 5.3, .24, "OBSERVED RESULT", 9, RED, True); text(slide, .6, 4.94, 5.35, .7, verdict, 16, WHITE, True)
    path = SHOTS / detail_file
    with Image.open(path) as image:
        iw, ih = image.size
    max_w, max_h = 6.05, 5.6; scale = min(max_w/iw, max_h/ih); width, height = iw*scale, ih*scale
    slide.shapes.add_picture(str(path), Inches(6.65+(max_w-width)/2), Inches(1.45+(max_h-height)/2), width=Inches(width), height=Inches(height))
    frame = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(6.58), Inches(1.37), Inches(6.2), Inches(5.75)); frame.fill.background(); frame.line.color.rgb = rgb(RED); frame.line.transparency = 35
    text(slide, .55, 7.08, 7.0, .16, "LIVE LOCAL UI CAPTURE  •  127.0.0.1:5174 → 127.0.0.1:5000", 7, MUTED, True)


def main():
    prs = Presentation(DECK)
    entries = [
        (24, "Dashboard overview — live frontend", "Opened the local React security dashboard at http://127.0.0.1:5174.", "01_dashboard.png", "The dashboard loaded successfully with sidebar navigation and live service status."),
        (25, "URL scan — safe example", "https://example.com", "02_url_safe_result_detail.png", "The report returned a Safe verdict with an explainable informational finding."),
        (26, "URL scan — suspicious pattern example", "http://paypal-security-check.example.invalid/login?verify=account", "03_url_high_risk_result_detail.png", "The local heuristics identified risky structural signals and produced a high-risk response."),
        (27, "Message scan — phishing-style test", "Urgent account-review text plus invoice.pdf.exe attachment metadata.", "04_message_result_detail.png", "Message, URL, and attachment indicators were combined into one defensive report."),
        (28, "File scan — safe static-analysis test", "tests/fixtures/safe_test.txt", "05_file_result_detail.png", "A harmless sample file was inspected statically; no file execution occurred."),
        (29, "QR scan — decoded text test", "Decoded QR payload: https://example.com", "06_qr_result_detail.png", "The decoded QR text was analyzed using the same safe URL-focused workflow."),
        (30, "Report lookup — persisted scan retrieval", "A scan ID captured from the URL test was submitted to the report lookup screen.", "07_report_lookup_result_detail.png", "The API returned the original saved report by ID, validating the report retrieval workflow."),
    ]
    for item in entries:
        slide = prs.slides.add_slide(prs.slide_layouts[6]); add(slide, *item)
    prs.save(DECK); print(DECK, len(prs.slides))


if __name__ == "__main__": main()
