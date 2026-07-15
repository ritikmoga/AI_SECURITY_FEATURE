"""Create the ScamShield AI project presentation (22 diagram-led slides)."""
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.oxml.xmlchemy import OxmlElement


OUT = Path(__file__).resolve().parents[1] / "docs" / "ScamShield_AI_Project_Presentation.pptx"
BANNER = Path(r"C:\Users\RITIK\Desktop\ritikmoga-profile\assets\cybersecurity-3d-banner.png")
W, H = 13.333, 7.5
BG, PANEL, RED, PINK, WHITE, MUTED = "080304", "1A070A", "E11D3A", "FF8A99", "FFF5F5", "BFA6A9"


def color(hex_value): return RGBColor.from_string(hex_value)


def text_box(slide, x, y, w, h, text, size=18, color_value=WHITE, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame; frame.clear(); frame.word_wrap = True
    p = frame.paragraphs[0]; p.alignment = align
    run = p.add_run(); run.text = text; run.font.name = "Aptos"; run.font.size = Pt(size); run.font.bold = bold; run.font.color.rgb = color(color_value)
    return box


def fill(shape, hex_value, transparency=0):
    shape.fill.solid(); shape.fill.fore_color.rgb = color(hex_value); shape.fill.transparency = transparency
    shape.line.color.rgb = color(RED); shape.line.transparency = 25


def add_transition(slide):
    transition = OxmlElement("p:transition"); transition.set("spd", "slow"); transition.set("advClick", "1")
    fade = OxmlElement("p:fade"); transition.append(fade)
    slide._element.insert(-1, transition)


def add_bg(slide, number):
    bg = slide.background.fill; bg.solid(); bg.fore_color.rgb = color(BG)
    accent = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(W), Inches(.08)); accent.fill.solid(); accent.fill.fore_color.rgb = color(RED); accent.line.fill.background()
    text_box(slide, 11.9, .22, 1, .25, f"{number:02d} / 22", 8, MUTED, True, PP_ALIGN.RIGHT)
    text_box(slide, .45, 7.08, 6, .18, "SCAMSHIELD AI  •  DEFENSIVE SECURITY PLATFORM", 7, MUTED, True)
    add_transition(slide)


def add_title(slide, title, kicker, number):
    add_bg(slide, number); text_box(slide, .55, .48, 11.5, .35, kicker.upper(), 10, RED, True)
    text_box(slide, .55, .83, 11.9, .7, title, 30, WHITE, True)


def add_bullets(slide, bullets):
    y = 1.72
    for item in bullets:
        text_box(slide, .75, y, 5.0, .5, "• " + item, 15, PINK); y += .62


def node(slide, x, y, label, sub="", width=1.65, height=.78, emphasis=False):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(width), Inches(height)); fill(shape, RED if emphasis else PANEL)
    if emphasis: shape.fill.transparency = 8
    text_box(slide, x+.08, y+.14, width-.16, .28, label, 12, WHITE, True, PP_ALIGN.CENTER)
    if sub: text_box(slide, x+.08, y+.44, width-.16, .18, sub, 7, MUTED, False, PP_ALIGN.CENTER)


def arrow(slide, x1, y1, x2, y2):
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2)); line.line.color.rgb = color(RED); line.line.width = Pt(1.8); line.line.end_arrowhead = True


def diagram(slide, labels, mode="flow"):
    # Connectors are created first, then nodes sit cleanly above them.
    if mode == "hub":
        cx, cy = 9.35, 4.05
        positions = [(7.15, 2.25), (9.35, 1.85), (11.55, 2.25), (7.15, 5.05), (9.35, 5.45), (11.55, 5.05)]
        for x, y in positions: arrow(slide, cx+.8, cy+.38, x+.8, y+.35)
        node(slide, cx, cy, labels[0], "core", 1.7, .8, True)
        for (x, y), label in zip(positions, labels[1:]): node(slide, x, y, label, "module", 1.6, .72)
    elif mode == "cycle":
        positions = [(8.2, 2.15), (10.65, 2.65), (10.65, 5.0), (8.2, 5.45), (6.55, 3.8)]
        for i, (x, y) in enumerate(positions):
            nx, ny = positions[(i+1) % len(positions)]; arrow(slide, x+.75, y+.35, nx+.75, ny+.35)
        for i, (x, y) in enumerate(positions): node(slide, x, y, labels[i], "step", 1.55, .72, i == 0)
    else:
        start_x, y = 6.25, 3.75; gap = 1.68
        for i in range(len(labels)-1): arrow(slide, start_x+i*gap+1.38, y+.36, start_x+(i+1)*gap, y+.36)
        for i, label in enumerate(labels): node(slide, start_x+i*gap, y, label, "stage", 1.35, .72, i == len(labels)-1)


def slide(prs, number, kicker, title, bullets, labels, mode="flow"):
    s = prs.slides.add_slide(prs.slide_layouts[6]); add_title(s, title, kicker, number); add_bullets(s, bullets); diagram(s, labels, mode); return s


def main():
    prs = Presentation(); prs.slide_width = Inches(W); prs.slide_height = Inches(H)
    # Cover
    cover = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(cover, 1)
    if BANNER.exists(): cover.shapes.add_picture(str(BANNER), Inches(0), Inches(.08), width=Inches(W), height=Inches(4.75))
    overlay = cover.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(.08), Inches(W), Inches(4.75)); overlay.fill.solid(); overlay.fill.fore_color.rgb = color(BG); overlay.fill.transparency = 35; overlay.line.fill.background()
    text_box(cover, .7, 4.95, 11.8, .72, "SCAMSHIELD AI", 38, WHITE, True)
    text_box(cover, .72, 5.7, 10.8, .36, "A defensive security-scanning platform for modern digital threats", 19, PINK)
    text_box(cover, .72, 6.34, 7, .3, "Project presentation  •  Ritik Moga", 12, MUTED, True)
    diagram(cover, ["SCAN", "ANALYZE", "SCORE", "PROTECT"], "flow")
    data = [
        ("01  •  ROADMAP", "What this deck covers", ["The security problem and product vision", "Platform architecture, scans, and risk scoring", "Authentication, persistence, deployment, and roadmap"], ["Problem", "Platform", "Security", "Scale"], "flow"),
        ("02  •  PROBLEM", "Digital threats move faster than human review", ["Links, messages, attachments, and QR codes can hide risk", "Users need a clear, fast, defensive first check", "The goal is triage—not executing suspicious content"], ["Input", "Unknown risk", "Human delay", "Exposure"], "flow"),
        ("03  •  VISION", "One defensive workspace for four common attack surfaces", ["A single console reduces switching cost", "Each input type receives purpose-built analysis", "Risk reports translate signals into practical next actions"], ["ScamShield", "URLs", "Messages", "Files", "QR codes", "Reports"], "hub"),
        ("04  •  ARCHITECTURE", "A layered platform designed for safe analysis", ["React + TypeScript delivers the operator experience", "Flask coordinates validation, analysis, scoring, and storage", "Optional providers extend local heuristics without becoming a dependency"], ["Browser", "API", "Validation", "Detectors", "Scoring", "Storage"], "flow"),
        ("05  •  FRONTEND", "A crimson security operations dashboard", ["Persistent sidebar creates dedicated scanner workspaces", "Results use risk badges, gauges, findings, and metadata", "Google sign-in enables personal history when configured"], ["Sidebar", "Scanner page", "API client", "Result view", "History"], "flow"),
        ("06  •  NAVIGATION", "Separate pages keep every scanning task focused", ["URL, message, file, and QR workflows are not hidden in one tab", "Reports and lookup are separate destinations", "Administrators receive a platform-posture destination"], ["Dashboard", "URL", "Message", "File", "QR", "Reports", "Admin"], "hub"),
        ("07  •  URL SCANNING", "URL intelligence combines syntax and safety checks", ["Validates scheme and host before analysis", "Flags suspicious patterns, redirect bait, and phishing signals", "Avoids unsafe fetching of local or private network targets"], ["URL", "Validate", "Heuristics", "Findings", "Risk report"], "flow"),
        ("08  •  MESSAGE SCANNING", "Message analysis surfaces social-engineering indicators", ["Inspects subject, body, attachment names, and embedded URLs", "Findings remain explainable instead of opaque", "Output recommends safe next actions"], ["Message", "Extract", "URL checks", "Malware cues", "Report"], "flow"),
        ("09  •  FILE SCANNING", "Static inspection protects the scanner and the user", ["Uploads are saved only to a temporary location", "The service never executes uploaded files", "Allowed extensions and size caps reduce attack surface"], ["Upload", "Allowlist", "Temp store", "Static inspect", "Delete"], "flow"),
        ("10  •  QR SCANNING", "QR images are decoded locally before server analysis", ["The browser decodes QR image content with jsQR", "Only decoded text is sent to the backend", "The resulting text is inspected like a URL or message"], ["QR image", "Browser decode", "Text", "Heuristics", "Report"], "flow"),
        ("11  •  RISK ENGINE", "Risk scoring turns many signals into one clear decision", ["Each finding contributes a weighted score", "Severity bonuses preserve serious indicators", "Risk bands keep the response understandable"], ["Findings", "Weights", "0–100 score", "Risk band", "Action"], "flow"),
        ("12  •  REPORTS", "Reports are useful security artifacts—not just a score", ["Every report has a scan ID and timestamp", "Findings retain evidence and severity", "Signed-in users can revisit stored scan history"], ["Scan", "JSON report", "Local archive", "User history", "Lookup"], "flow"),
        ("13  •  AUTHENTICATION", "Google Sign-In is verified server-side", ["Browser receives a Google ID credential", "Flask verifies the ID token audience and email", "The app issues short-lived access and refresh JWTs"], ["Google", "ID token", "Flask verify", "JWT pair", "Session"], "flow"),
        ("14  •  DATA", "Storage scales from local development to PostgreSQL", ["SQLite gives the project a zero-setup local path", "PostgreSQL schema models users, reports, sessions, keys, and audit events", "Migration adds Google identity fields for production"], ["SQLite", "User store", "Reports", "PostgreSQL", "Audit log"], "flow"),
        ("15  •  RATE LIMITING", "Abuse protection has a development and production path", ["In-memory windows keep local demos safe", "Redis enables distributed limits across API instances", "Fallback behavior keeps the service available during Redis outages"], ["Client", "Limiter", "Redis", "Allow", "429 retry"], "flow"),
        ("16  •  API SECURITY", "Defensive defaults are built into the request lifecycle", ["Input validation constrains size, shape, and allowed file types", "Request IDs improve traceability", "Security headers, CORS controls, and safe error messages reduce exposure"], ["Request", "Validate", "Trace ID", "Secure headers", "Response"], "flow"),
        ("17  •  DEVOPS", "Quality gates protect the project as it evolves", ["GitHub Actions runs backend tests and frontend builds", "Dependabot tracks dependency updates", "Pre-commit rules catch common repository mistakes"], ["Commit", "Pre-commit", "CI", "Build", "Deploy"], "cycle"),
        ("18  •  DEPLOYMENT", "Container-ready today, production-ready by configuration", ["Docker packages the Flask service", "Persistent data volume supports the local store", "Production uses HTTPS, strict CORS, managed Postgres, Redis, and secret management"], ["Source", "Docker", "API container", "Data", "Reverse proxy"], "flow"),
        ("19  •  OPENAPI", "A documented API improves integration and testing", ["OpenAPI describes health, scan, auth, refresh, and dashboard routes", "Typed frontend helpers centralize request parsing", "Consistent envelopes simplify client error handling"], ["OpenAPI", "Endpoints", "Typed client", "Frontend", "Consumers"], "flow"),
        ("20  •  THREAT MODEL", "Security decisions are driven by explicit assets and controls", ["Uploads: static-only handling and deletion", "Accounts: verified Google tokens and JWT expiry", "Reports: authenticated history, retention guidance, and audit-ready schema"], ["Assets", "Threats", "Controls", "Monitoring", "Review"], "cycle"),
        ("21  •  ROADMAP", "From capable project to production service", ["Connect real Google credentials and configure secrets", "Activate PostgreSQL and Redis infrastructure", "Add queue workers, isolated malware engines, and observability"], ["Credentials", "Infrastructure", "Workers", "Telemetry", "Scale"], "flow"),
    ]
    for idx, (kicker, title, bullets, labels, mode) in enumerate(data, start=2): slide(prs, idx, kicker, title, bullets, labels, mode)
    end = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(end, 22)
    text_box(end, .7, 1.0, 12, .6, "Build defensively. Explain clearly. Improve continuously.", 31, WHITE, True, PP_ALIGN.CENTER)
    text_box(end, 2.1, 1.9, 9.2, .42, "ScamShield AI is a project foundation for safer digital decisions.", 18, PINK, False, PP_ALIGN.CENTER)
    diagram(end, ["SCAN", "ANALYZE", "DECIDE", "PROTECT"], "flow")
    text_box(end, 3.35, 5.15, 6.7, .35, "github.com/ritikmoga/AI_SECURITY_FEATURE", 15, RED, True, PP_ALIGN.CENTER)
    OUT.parent.mkdir(parents=True, exist_ok=True); prs.save(OUT); print(OUT)


if __name__ == "__main__": main()
